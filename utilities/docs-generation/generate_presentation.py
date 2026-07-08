#!/usr/bin/env python3
"""
Business-Focused PowerPoint Presentation Generator
For non-technical stakeholders - focuses on features, benefits, and business value
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

class BusinessPowerPointGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Professional color scheme
        self.primary_color = RGBColor(25, 77, 137)      # Dark Blue
        self.accent_color = RGBColor(68, 144, 216)      # Light Blue
        self.success_color = RGBColor(51, 163, 102)     # Green
        self.warning_color = RGBColor(237, 125, 49)     # Orange
        self.text_color = RGBColor(51, 51, 51)          # Dark Gray
        self.light_text = RGBColor(255, 255, 255)       # White
        
    def add_title_slide(self, title, subtitle, company=""):
        """Add professional title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Add gradient effect with colored shapes
        bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.primary_color
        bg_shape.line.fill.background()
        
        # Add accent bar
        accent = slide.shapes.add_shape(1, Inches(0), Inches(5.5), Inches(10), Inches(2))
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.accent_color
        accent.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(60)
        title_p.font.bold = True
        title_p.font.color.rgb = self.light_text
        title_p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1.2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = subtitle
        subtitle_p.font.size = Pt(28)
        subtitle_p.font.color.rgb = self.text_color
        subtitle_p.alignment = PP_ALIGN.CENTER
        
        # Company/Date
        if company:
            footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.5))
            footer_frame = footer_box.text_frame
            footer_p = footer_frame.paragraphs[0]
            footer_p.text = f"{company} • February 16, 2026"
            footer_p.font.size = Pt(14)
            footer_p.font.color.rgb = self.light_text
            footer_p.alignment = PP_ALIGN.CENTER
        
    def add_feature_slide(self, feature_name, goal, problem, business_value, use_case):
        """Add feature slide with business focus"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Add background
        bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
        bg_shape.line.fill.background()
        
        # Feature name header
        header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.primary_color
        header_shape.line.fill.background()
        
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
        header_frame = header_box.text_frame
        header_p = header_frame.paragraphs[0]
        header_p.text = f"🎯 {feature_name}"
        header_p.font.size = Pt(36)
        header_p.font.bold = True
        header_p.font.color.rgb = self.light_text
        
        # Content sections
        y_pos = 1.2
        section_height = 0.35
        content_height = 0.9
        
        # Goal
        self._add_section(slide, "GOAL", goal, y_pos, self.accent_color)
        y_pos += section_height + content_height
        
        # Problem Solved
        self._add_section(slide, "PROBLEM SOLVED", problem, y_pos, self.warning_color)
        y_pos += section_height + content_height
        
        # Business Value
        self._add_section(slide, "BUSINESS VALUE", business_value, y_pos, self.success_color)
        y_pos += section_height + content_height
        
        # Use Case
        self._add_section(slide, "USE CASE", use_case, y_pos, self.primary_color)
    
    def _add_section(self, slide, section_title, content, y_start, color):
        """Helper to add content section"""
        # Section title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_start), Inches(2), Inches(0.3))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = section_title
        title_p.font.size = Pt(12)
        title_p.font.bold = True
        title_p.font.color.rgb = color
        
        # Content box with background
        content_shape = slide.shapes.add_shape(1, Inches(0.7), Inches(y_start + 0.35), Inches(8.6), Inches(0.85))
        content_shape.fill.solid()
        content_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        content_shape.line.color.rgb = RGBColor(200, 200, 200)
        content_shape.line.width = Pt(1)
        
        # Content text
        content_box = slide.shapes.add_textbox(Inches(0.9), Inches(y_start + 0.45), Inches(8.2), Inches(0.7))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_p = content_frame.paragraphs[0]
        content_p.text = content
        content_p.font.size = Pt(11)
        content_p.font.color.rgb = self.text_color
    
    def add_ai_advantages_slide(self):
        """Add AI advantages slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
        bg_shape.line.fill.background()
        
        # Header
        header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.primary_color
        header_shape.line.fill.background()
        
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        header_frame = header_box.text_frame
        header_p = header_frame.paragraphs[0]
        header_p.text = "🤖 AI-Powered Intelligence Advantage"
        header_p.font.size = Pt(40)
        header_p.font.bold = True
        header_p.font.color.rgb = self.light_text
        
        # Advantages
        advantages = [
            ("🎯 Real-Time Analysis", "Processes information 24/7 automatically"),
            ("⚡ Speed & Accuracy", "AI analyzes thousands of sources instantly"),
            ("🧠 Smart Insights", "Machine learning identifies patterns humans miss"),
            ("💰 Cost Reduction", "Eliminates manual research and data entry"),
            ("📊 Data-Driven Decisions", "Supports decisions with AI-verified insights"),
            ("🔀 Competitor Sensing", "Automatic competitor mention detection")
        ]
        
        x_positions = [0.7, 5.2]
        y_start = 1.2
        
        for idx, (title, desc) in enumerate(advantages):
            col = idx % 2
            row = idx // 2
            
            x = x_positions[col]
            y = y_start + (row * 1.85)
            
            # Box
            box_shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(4), Inches(1.65))
            box_shape.fill.solid()
            box_shape.fill.fore_color.rgb = self.accent_color
            box_shape.line.color.rgb = self.primary_color
            box_shape.line.width = Pt(2)
            
            # Title
            title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(3.6), Inches(0.35))
            title_frame = title_box.text_frame
            title_p = title_frame.paragraphs[0]
            title_p.text = title
            title_p.font.size = Pt(13)
            title_p.font.bold = True
            title_p.font.color.rgb = self.light_text
            
            # Description
            desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.55), Inches(3.6), Inches(0.95))
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            desc_p = desc_frame.paragraphs[0]
            desc_p.text = desc
            desc_p.font.size = Pt(10)
            desc_p.font.color.rgb = self.light_text
    
    def add_roi_slide(self):
        """Add ROI and business impact slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
        bg_shape.line.fill.background()
        
        # Header
        header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.primary_color
        header_shape.line.fill.background()
        
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        header_frame = header_box.text_frame
        header_p = header_frame.paragraphs[0]
        header_p.text = "💼 Business Impact & ROI"
        header_p.font.size = Pt(40)
        header_p.font.bold = True
        header_p.font.color.rgb = self.light_text
        
        # ROI metrics
        metrics = [
            ("70%", "Time Saved", "Manual research automation", self.success_color),
            ("3x", "Faster Decision", "Real-time intelligence", self.warning_color),
            ("50+", "Competitors Tracked", "Automated monitoring", self.accent_color),
            ("24/7", "Intelligence Coverage", "Non-stop analysis", RGBColor(155, 89, 182))
        ]
        
        for idx, (value, metric, desc, color) in enumerate(metrics):
            col = idx % 2
            row = idx // 2
            
            x = 0.8 + (col * 4.8)
            y = 1.3 + (row * 2.8)
            
            # Metric box
            metric_shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.2), Inches(2.3))
            metric_shape.fill.solid()
            metric_shape.fill.fore_color.rgb = color
            metric_shape.line.fill.background()
            
            # Value
            value_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.25), Inches(3.6), Inches(0.7))
            value_frame = value_box.text_frame
            value_p = value_frame.paragraphs[0]
            value_p.text = value
            value_p.font.size = Pt(48)
            value_p.font.bold = True
            value_p.font.color.rgb = self.light_text
            value_p.alignment = PP_ALIGN.CENTER
            
            # Metric name
            metric_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.9), Inches(3.6), Inches(0.4))
            metric_frame = metric_box.text_frame
            metric_p = metric_frame.paragraphs[0]
            metric_p.text = metric
            metric_p.font.size = Pt(14)
            metric_p.font.bold = True
            metric_p.font.color.rgb = self.light_text
            metric_p.alignment = PP_ALIGN.CENTER
            
            # Description
            desc_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 1.35), Inches(3.6), Inches(0.8))
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            desc_p = desc_frame.paragraphs[0]
            desc_p.text = desc
            desc_p.font.size = Pt(9)
            desc_p.font.color.rgb = self.light_text
            desc_p.alignment = PP_ALIGN.CENTER
    
    def add_technical_slide(self):
        """Add 1 technical overview slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
        bg_shape.line.fill.background()
        
        # Header
        header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.primary_color
        header_shape.line.fill.background()
        
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        header_frame = header_box.text_frame
        header_p = header_frame.paragraphs[0]
        header_p.text = "⚙️ Technical Stack (Overview)"
        header_p.font.size = Pt(40)
        header_p.font.bold = True
        header_p.font.color.rgb = self.light_text
        
        # Left column - Architecture
        arch_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(4.5), Inches(5.8))
        arch_frame = arch_box.text_frame
        arch_frame.word_wrap = True
        
        sections = [
            ("Backend", ".NET 8 / ASP.NET Core"),
            ("Frontend", "Angular 17 (Modern TypeScript)"),
            ("Database", "SQL Server (Structured Data)"),
            ("AI Engine", "Google Gemini API (Analysis)"),
            ("Cloud", "Azure Blob Storage (Files)"),
            ("Real-time", "SignalR WebSockets"),
            ("Automation", "Python Workers (24/7)")
        ]
        
        for idx, (name, tech) in enumerate(sections):
            if idx == 0:
                p = arch_frame.paragraphs[0]
            else:
                p = arch_frame.add_paragraph()
            
            p.text = f"• {name}: {tech}"
            p.font.size = Pt(11)
            p.font.color.rgb = self.text_color
            p.space_after = Pt(10)
        
        # Right column - Data Flow
        flow_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.2), Inches(4), Inches(0.4))
        flow_title_frame = flow_title_box.text_frame
        flow_title_p = flow_title_frame.paragraphs[0]
        flow_title_p.text = "Data Flow & Updates"
        flow_title_p.font.size = Pt(13)
        flow_title_p.font.bold = True
        flow_title_p.font.color.rgb = self.primary_color
        
        flow_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.7), Inches(4), Inches(5.3))
        flow_frame = flow_box.text_frame
        flow_frame.word_wrap = True
        
        flow_items = [
            "1. User adds company to track",
            "2. Python worker fetches live data",
            "3. AI analyzes content & trends",
            "4. Sentiment analyzed in real-time",
            "5. Alerts triggered for key events",
            "6. Reports generated automatically",
            "7. PDFs stored in cloud",
            "8. Dashboard updates live",
            "",
            "Updates: Every 5 minutes (configurable)"
        ]
        
        for idx, item in enumerate(flow_items):
            if idx == 0:
                p = flow_frame.paragraphs[0]
            else:
                p = flow_frame.add_paragraph()
            
            p.text = item
            p.font.size = Pt(10)
            p.font.color.rgb = self.text_color
            p.space_after = Pt(6)
    
    def add_closing_slide(self):
        """Add closing slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background gradient effect
        bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.primary_color
        bg_shape.line.fill.background()
        
        # Main message
        main_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        main_frame = main_box.text_frame
        main_frame.word_wrap = True
        main_p = main_frame.paragraphs[0]
        main_p.text = "Transform Your Market Intelligence"
        main_p.font.size = Pt(48)
        main_p.font.bold = True
        main_p.font.color.rgb = self.light_text
        main_p.alignment = PP_ALIGN.CENTER
        
        # Subtext
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.6), Inches(8), Inches(1.5))
        sub_frame = sub_box.text_frame
        sub_frame.word_wrap = True
        sub_p = sub_frame.paragraphs[0]
        sub_p.text = "Make faster, smarter decisions with AI-powered market insights\n\nAlfanar Market Intelligence • Your Competitive Edge"
        sub_p.font.size = Pt(18)
        sub_p.font.color.rgb = self.light_text
        sub_p.alignment = PP_ALIGN.CENTER
    
    def generate(self):
        """Generate complete business presentation"""
        
        # Slide 1: Title
        self.add_title_slide(
            "Alfanar Market Intelligence",
            "AI-Powered Business Intelligence Platform",
            "Alfanar Tech Solutions"
        )
        
        # Slide 2: Overview
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
        bg_shape.line.fill.background()
        
        header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.primary_color
        header_shape.line.fill.background()
        
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        header_frame = header_box.text_frame
        header_p = header_frame.paragraphs[0]
        header_p.text = "📱 Platform Overview"
        header_p.font.size = Pt(40)
        header_p.font.bold = True
        header_p.font.color.rgb = self.light_text
        
        overview_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(5.8))
        overview_frame = overview_box.text_frame
        overview_frame.word_wrap = True
        
        overview_items = [
            "Every day, your competitors are making strategic moves.",
            "New market trends emerge. Acquisition opportunities appear. Risks develop.",
            "",
            "Alfanar Market Intelligence gives you a 360° view of your market landscape.",
            "",
            "Our AI continuously monitors competitors, analyzes market trends,",
            "detects opportunities, and alerts you to action items—all in real-time.",
            "",
            "Transform raw market data into actionable business intelligence."
        ]
        
        for idx, item in enumerate(overview_items):
            if idx == 0:
                p = overview_frame.paragraphs[0]
            else:
                p = overview_frame.add_paragraph()
            
            if item == "":
                p.text = " "
                p.space_after = Pt(8)
            else:
                p.text = item
                p.font.size = Pt(14)
                p.font.color.rgb = self.text_color
                p.space_after = Pt(10)
        
        # Slide 3: Feature 1 - Intelligence Reports
        self.add_feature_slide(
            "Intelligence Reports",
            "Generate comprehensive market analysis reports on-demand",
            "Manual research takes weeks. You need insights NOW.",
            "Instant market analysis + AI-powered insights = Competitive advantage",
            "Marketing team needs STATCOM market overview for Q1 strategy. System generates 5-page report in 2 minutes."
        )
        
        # Slide 4: Feature 2 - Competitor Tracking
        self.add_feature_slide(
            "Competitor Tracking",
            "Monitor competitor activities, mentions, and sentiment automatically",
            "You hear about competitor moves weeks late via industry news.",
            "Real-time competitor intelligence enables faster counter-moves",
            "ABB announces new STATCOM product. System alerts you within 1 hour. You brief executives before market reacts."
        )
        
        # Slide 5: Feature 3 - Trend Analysis
        self.add_feature_slide(
            "Market Trend Analysis",
            "Identify emerging market trends before competitors",
            "Trend reports from consultants arrive too late. You're always reacting.",
            "Early trend detection = First-mover advantage in market shifts",
            "Detect rising demand for renewable energy STATCOM solutions 6 months before industry reports. Launch product early."
        )
        
        # Slide 6: Feature 4 - Technology Intelligence
        self.add_feature_slide(
            "Technology Intelligence",
            "Track technological advancements and innovation patterns",
            "Your team misses emerging tech disruptions until it's too late.",
            "Stay ahead of tech disruption with early pattern recognition",
            "AI detects breakthrough in battery storage tech. R&D team investigates and forms partnership before competitors notice."
        )
        
        # Slide 7: Feature 5 - Keyword Monitoring
        self.add_feature_slide(
            "Keyword & Web Monitoring",
            "Track specific keywords and topics across the entire web",
            "You can't manually track thousands of articles, news, and mentions.",
            "Real-time web monitoring = Complete market visibility",
            "Monitor 'grid modernization' across 100+ sources. Get weekly summaries with key developments."
        )
        
        # Slide 8: Feature 6 - Automated Alerts
        self.add_feature_slide(
            "Automated Intelligence Alerts",
            "Get notified of critical market events and sentiment shifts",
            "Important news gets buried in email. You miss critical developments.",
            "Don't miss market-moving events. React in real-time.",
            "Stock price drops dramatically. AI correlates to negative sentiment spike about customer. Alert goes to executives in seconds."
        )
        
        # Slide 9: PowerPoint Feature
        self.add_feature_slide(
            "PowerPoint Presentation Generation",
            "Export reports to professional PowerPoint presentations for stakeholders",
            "Creating executive presentations takes hours. You're copying-pasting data manually.",
            "Turn data into presentation-ready documents in seconds",
            "Board meeting tomorrow. Generate STATCOM market analysis in PowerPoint in 2 minutes with charts, competitor rankings, and recommendations."
        )
        
        # Slide 10: AI Advantages
        self.add_ai_advantages_slide()
        
        # Slide 11: Business Impact & ROI
        self.add_roi_slide()
        
        # Slide 12: Tech Stack (1 slide only)
        self.add_technical_slide()
        
        # Slide 13: Closing
        self.add_closing_slide()
        
        # Save presentation
        output_path = r"D:\Storage Market Intel\Alfanar.MarketIntel\POWERPOINT_FEATURE_PRESENTATION.pptx"
        self.prs.save(output_path)
        print(f"✅ Business PowerPoint presentation created successfully!")
        print(f"📁 Location: {output_path}")
        print(f"📊 Total slides: {len(self.prs.slides)}")

def main():
    print("🎬 Generating business-focused PowerPoint presentation...")
    generator = BusinessPowerPointGenerator()
    generator.generate()

if __name__ == "__main__":
    main()
